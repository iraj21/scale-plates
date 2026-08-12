#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Build the Scale Plates client pitch deck.

Recreated from the original Atlas deck's design system (geometry, palette,
typography) — same story arc, renamed to SCALE PLATES, Zomato + Swiggy
throughout, and no methodology/benchmark numbers (restaurant counts, months of
data) anywhere. Case studies are anonymized (Restaurant A / B).

Output: pitch/Scale-Plates-Client-Pitch.pptx (16:9, 11 slides)
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- palette
NAVY_D = RGBColor(0x14, 0x33, 0x52)   # slide bg / deep
NAVY = RGBColor(0x1F, 0x4E, 0x79)     # panels / accents
NAVY2 = RGBColor(0x2A, 0x5A, 0x8C)    # lighter navy panel
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
INK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREY2 = RGBColor(0x8F, 0xA6, 0xC4)
BLUE1 = RGBColor(0xC9, 0xD6, 0xE8)
BLUE2 = RGBColor(0xDD, 0xE7, 0xF2)
CARD = RGBColor(0xF4, 0xF6, 0xF8)
TRACK_BG = RGBColor(0xE3, 0xE8, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(para, str):
            para = [(para, {})]
        for txt, st in para:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = st.get("font", FONT)
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", bold)
            f.color.rgb = st.get("color", color)
    return tb


def dot(s, x, y, color, d=0.13):
    rect(s, x, y, d, d, color)


# ---------------------------------------------------------------- shared chrome
def header(s, kicker, title, num):
    rect(s, 0, 0, 13.333, 0.95, WHITE)
    rect(s, 0, 0, 13.333, 0.06, ORANGE)
    text(s, 0.55, 0.12, 9.5, 0.3, kicker, size=11, color=ORANGE, bold=True)
    text(s, 0.55, 0.35, 10.5, 0.6, title, size=25, color=NAVY, bold=True)
    text(s, 12.1, 0.32, 0.9, 0.4, num, size=11, color=GREY, align=PP_ALIGN.RIGHT)
    text(s, 0.55, 6.95, 12.2, 0.35,
         "SCALE PLATES  •  We help restaurants keep more money  •  Same kitchen, more profit",
         size=9, color=GREY)


def strip_card(s, x, y, w, h, strip, fill=CARD):
    rect(s, x, y, w, h, fill)
    rect(s, x, y, w, 0.09, strip)
    rect(s, x, y, 0.09, h, strip)  # left edge accent (matches Atlas cards)


def big_panel(s, x, y, w, h, fill):
    rect(s, x, y, w, h, fill)


# ================================================================ 01 TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY_D)
rect(s, 0, 0, 13.333, 0.12, ORANGE)
rect(s, 10.0, 0, 3.33, 7.5, NAVY)
_LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "pitch", "logo_transparent.png")
if os.path.exists(_LOGO):
    # white rounded badge so the dark logo reads on the navy background
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(1.75), Inches(1.75))
    badge.adjustments[0] = 0.12
    badge.fill.solid()
    badge.fill.fore_color.rgb = WHITE
    badge.line.fill.background()
    badge.shadow.inherit = False
    s.shapes.add_picture(_LOGO, Inches(0.92), Inches(0.57), Inches(1.5), Inches(1.5))
else:
    text(s, 0.9, 0.75, 6.0, 0.5, "SCALE PLATES", size=30, color=WHITE, bold=True)
text(s, 0.92, 2.45, 6.0, 0.3, "WE HELP RESTAURANTS MAKE MORE MONEY", size=11, color=ORANGE, bold=True)
text(s, 0.9, 2.95, 8.8, 1.6, [
    [("More money.", {"size": 52, "color": WHITE, "bold": True})],
    [("Same kitchen.", {"size": 52, "color": GOLD, "bold": True})],
], spacing=1.0)
text(s, 0.9, 4.35, 8.2, 0.9,
     "We read your Zomato and Swiggy reports, and show you where your money is going — "
     "and how to keep more of it. Your food does not change.", size=16, color=BLUE1)
big_panel(s, 10.55, 2.0, 2.3, 3.4, NAVY2)
text(s, 10.85, 2.35, 1.7, 0.4, "Our promise", size=13, color=ORANGE, bold=True)
text(s, 10.85, 2.85, 1.75, 2.3,
     "Every rupee that goes out — ads, offers, prices, fees — is decided with data and "
     "checked against profit.", size=13, color=WHITE)
text(s, 0.9, 6.7, 11.0, 0.35, "WE LOOK AFTER THE MONEY SIDE OF YOUR ONLINE RESTAURANT BUSINESS",
     size=10, color=GREY2, bold=True)

# ================================================================ 02 PROBLEM
s = slide()
header(s, "THE PROBLEM", "Why most restaurants lose money quietly", "02")
cards = [
    (0.55, RED, "~30%", "of every rupee you sell goes to the platform",
     "commission + payment fee + long-distance fee + tax on fees"),
    (4.71, ORANGE, "40%", "of orders pay a hidden ₹15–25 fee",
     "the customer never sees it, but it is cut from YOUR money"),
    (8.87, NAVY, "8 in 10", "orders now come from paid ads",
     "no organic backup — pause the ads and revenue pauses"),
]
for x, col, big, t, sub in cards:
    strip_card(s, x, 1.5, 3.95, 3.6, col)
    text(s, x + 0.3, 1.95, 3.4, 1.0, big, size=44, color=col, bold=True)
    text(s, x + 0.3, 3.0, 3.35, 0.9, t, size=15, color=INK, bold=True)
    text(s, x + 0.3, 3.9, 3.35, 1.1, sub, size=11.5, color=GREY)
text(s, 0.55, 5.5, 12.2, 0.7,
     "You are a great chef. But who watches your prices, offers, ads, fees and orders? "
     "That is the work we do for you.", size=15, color=NAVY, bold=True)
text(s, 0.55, 6.62, 12.2, 0.3,
     "Your food is fine. The problem: nobody is watching your online business like a business.",
     size=10.5, color=GREY)

# ================================================================ 03 SOLUTION
s = slide()
header(s, "THE SOLUTION", "Ads are the engine. The rest is the body of the car.", "03")
big_panel(s, 0.55, 1.5, 5.9, 4.0, NAVY)
text(s, 0.85, 1.8, 5.3, 0.4, "THE ENGINE — YOUR AD MONEY", size=13, color=ORANGE, bold=True)
text(s, 0.85, 2.25, 5.3, 0.6, "We spend your ad money carefully — only where it brings orders.",
     size=15, color=WHITE, bold=True)
bullets = [
    "Every ad checked: earn below ₹3 per ₹1 → stop, above ₹4 → increase",
    "Put 70% of ad money in your peak window (dinner or lunch — from your data)",
    "Reported returns can claim 8–9× — we count only the EXTRA orders ads truly bring",
    "We stop the ad before it stops giving returns",
]
y = 3.06
for b in bullets:
    dot(s, 0.85, y + 0.02, GOLD)
    text(s, 1.15, y - 0.06, 5.1, 0.55, b, size=12, color=BLUE2)
    y += 0.62
big_panel(s, 6.7, 1.5, 6.08, 4.0, CARD)
text(s, 7.0, 1.8, 5.5, 0.4, "THE BODY — WHAT MAKES ADS WORK", size=13, color=ORANGE, bold=True)
body = [
    ("Menu & price", "combos and offers, so ad customers actually order"),
    ("Offers", "use the platform's money for offers, never your profit"),
    ("Fees & wastage", "stop hidden fees and wrong packing charges"),
    ("Order flow", "turn more viewers into paying customers"),
    ("Daily work", "fast replies, good ratings, on-time food"),
]
y = 2.3
for lab, desc in body:
    text(s, 7.0, y, 1.9, 0.4, lab, size=12.5, color=NAVY, bold=True)
    text(s, 8.9, y, 3.7, 0.55, desc, size=11, color=GREY)
    y += 0.64
big_panel(s, 0.55, 5.7, 12.23, 0.85, ORANGE)
text(s, 0.85, 5.88, 11.7, 0.5,
     "Ads bring more orders. We make sure you also keep more profit. Every ad rupee is "
     "checked against money in your bank.", size=14, color=WHITE, bold=True)
text(s, 0.55, 6.62, 12.2, 0.3,
     "Ad agencies sell views. We increase the money you actually receive.", size=10.5, color=GREY)

# ================================================================ 04 OFFER / TRACKS
s = slide()
header(s, "THE OFFER", "Two ways to work with us — choose your goal", "04")
track1 = [
    "Use data to cut or re-time ad spend",
    "Stop hidden long-distance fees",
    "Increase order size with combos and add-ons",
    "Fix offers and packing charges",
    "More viewers become paying customers",
]
strip_card(s, 0.55, 1.5, 6.0, 4.9, NAVY)
text(s, 0.85, 1.75, 5.4, 0.4, "TRACK 1", size=12, color=NAVY, bold=True)
text(s, 0.85, 2.1, 5.4, 0.6, "MAKE MORE PROFIT", size=20, color=NAVY, bold=True)
text(s, 0.85, 2.75, 5.4, 0.4, "Same orders, more money in your pocket.", size=12.5, color=GREY)
y = 3.25
for b in track1:
    dot(s, 0.87, y + 0.05, NAVY)
    text(s, 1.15, y, 5.2, 0.35, b, size=12, color=INK)
    y += 0.42
rect(s, 0.85, 5.35, 5.4, 0.85, CARD)
text(s, 1.05, 5.5, 5.0, 0.6, "For: a running business that wants more profit from the same orders",
     size=11, color=GREY)
track2 = [
    "Increase ads where they already work",
    "Offers that pay for themselves",
    "More customers, and regulars coming back",
    "Expand to new areas",
    "More viewers become paying customers",
]
strip_card(s, 6.75, 1.5, 6.0, 4.9, ORANGE)
text(s, 7.05, 1.75, 5.4, 0.4, "TRACK 2", size=12, color=ORANGE, bold=True)
text(s, 7.05, 2.1, 5.4, 0.6, "GROW MORE", size=20, color=NAVY, bold=True)
text(s, 7.05, 2.75, 5.4, 0.4, "Spend more on ads and offers — but only where they work.",
     size=12.5, color=GREY)
y = 3.25
for b in track2:
    dot(s, 7.07, y + 0.05, ORANGE)
    text(s, 7.35, y, 5.2, 0.35, b, size=12, color=INK)
    y += 0.42
rect(s, 7.05, 5.35, 5.4, 0.85, CARD)
text(s, 7.25, 5.5, 5.0, 0.6, "For: ready to invest and grow orders and area", size=11, color=GREY)
text(s, 0.55, 6.62, 12.2, 0.3,
     "Both start with the same free check-up of your business. You choose the goal, we make the plan.",
     size=10.5, color=GREY)

# ================================================================ 05 THE DATA / SCORING
s = slide()
header(s, "THE DATA", "Your payout report has all the answers. We find them.", "05")
big_panel(s, 0.55, 1.5, 5.6, 4.9, CARD)
text(s, 0.85, 1.8, 5.0, 0.4, "FROM YOUR ZOMATO / SWIGGY REPORT", size=12, color=ORANGE, bold=True)
data_pts = [
    "Every order has 60+ details — tens of thousands of data points every month",
    "We track every rupee: order value, fees, tax, distance, time, offers",
    "Ad invoices: what you spent, what you got back",
    "Weekly payments: what is paid, what is pending",
    "Every offer running, and who is paying for it",
]
y = 2.37
for b in data_pts:
    dot(s, 0.9, y + 0.05, ORANGE)
    text(s, 1.2, y, 4.8, 0.7, b, size=12.5, color=INK)
    y += 0.78
text(s, 0.85, 5.6, 5.0, 0.7,
     "We also use views, ratings, reviews and your billing data where available.", size=11, color=GREY)
big_panel(s, 6.5, 1.5, 6.3, 4.9, CARD)
text(s, 6.85, 1.8, 5.6, 0.4, "YOU GET A HEALTH SCORE (0–100)", size=12, color=ORANGE, bold=True)
text(s, 6.85, 2.2, 2.2, 1.0, "86", size=54, color=NAVY, bold=True)
text(s, 8.9, 2.5, 3.6, 0.5, "out of 100\nYour overall health", size=12, color=GREY)
dials = [
    ("Revenue", 98, GREEN), ("Ads", 68, ORANGE), ("Pricing", 88, GREEN),
    ("Menu / Radius", 100, GREEN), ("Operations", 98, GREEN), ("Profitability", 100, GREEN),
    ("Repeat", 100, GREEN), ("Rating", 92, GREEN),
]
y = 3.25
for lab, val, col in dials:
    text(s, 6.85, y, 2.4, 0.3, lab, size=11, color=INK, bold=True)
    rect(s, 9.0, y + 0.03, 3.3, 0.22, TRACK_BG)
    rect(s, 9.0, y + 0.03, 3.3 * val / 100, 0.22, col)
    text(s, 12.35, y, 0.45, 0.3, str(val), size=11, color=INK, align=PP_ALIGN.RIGHT)
    y += 0.38
text(s, 6.85, y + 0.02, 5.6, 0.6, "Red dial = money going out. That is where we start.",
     size=12, color=RED, bold=True)
text(s, 0.55, 6.62, 12.2, 0.3,
     "These scores are from a real restaurant's data — see Slide 10.", size=10.5, color=GREY)

# ================================================================ 06 THE FUNNEL
s = slide()
header(s, "THE FUNNEL", "Where customers walk away — and how we bring them back", "06")
funnel = [
    ("SEEN", "customers see you in search", "everyone", 2.9, "0.75"),
    ("OPEN", "they open your page", "~8%", 2.55, "0.2"),
    ("CART", "they add food", "30%+", 2.2, "0.2"),
    ("ORDER", "they pay", "55%+", 1.85, "0.1"),
]
x = 0.55
for i, (name, desc, val, w, pad) in enumerate(funnel):
    if i == 3:
        fill = ORANGE
    else:
        fill = NAVY
    big_panel(s, x, 1.9, w, 1.15, fill)
    text(s, x + 0.2, 2.05, w - 0.4, 0.35, name, size=13.5, color=WHITE, bold=True)
    text(s, x + 0.2, 2.42, w - 0.4, 0.3, desc, size=10, color=BLUE2)
    text(s, x + 0.2, 2.72, w - 0.4, 0.3, val, size=10.5, color=GOLD, bold=True)
    nx = x + w + 0.03
    if i < 3:
        big_panel(s, nx, 2.35, 0.45, 0.3, ORANGE)
        text(s, nx, 2.35, 0.45, 0.3, "→", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x = nx + 0.45
text(s, 0.55, 3.35, 12.2, 0.4, "These two numbers decide your growth. Both can be fixed:",
     size=13, color=NAVY, bold=True)
strip_card(s, 0.55, 3.9, 6.0, 2.4, ORANGE)
text(s, 0.85, 4.15, 5.4, 0.5, "If less than 30% of viewers order — menu page problem",
     size=14.5, color=ORANGE, bold=True)
text(s, 0.85, 4.7, 5.4, 1.5,
     "Better photos, clear dish names, best dishes on top, combos easy to see. We fix the menu page.",
     size=12, color=INK)
strip_card(s, 6.75, 3.9, 6.0, 2.4, RED)
text(s, 7.05, 4.15, 5.4, 0.5, "If less than 55% of carts are paid — checkout problem",
     size=14.5, color=RED, bold=True)
text(s, 7.05, 4.7, 5.4, 1.5,
     "Delivery fee shock, hidden charges at payment. We remove the reasons people walk away.",
     size=12, color=INK)
text(s, 0.55, 6.62, 12.2, 0.3,
     "More orders from the same customers = free growth. Most of the money is hiding here.",
     size=10.5, color=GREY)

# ================================================================ 07 THE AD MACHINE
s = slide()
header(s, "THE AD MACHINE", "Ad money: spent carefully, checked daily", "07")
rules = [
    (0.55, NAVY, "4×", "return rule", "Ads must return at least ₹4 for every ₹1. Below ₹3 → stop within 7 days."),
    (3.7, ORANGE, "70%", "peak window rule", "Put 70% of ad money in your peak window (dinner or lunch — from your data)."),
    (6.85, RED, "STOP", "check rule", "No ad runs more than 7 days without a check. We watch the numbers, we don't hope."),
    (10.0, GREEN, "10%", "test rule", "Test new dishes or offers with a small 10% slice of budget — scale only what works."),
]
for x, col, big, t, sub in rules:
    strip_card(s, x, 1.5, 2.9, 2.5, col)
    text(s, x + 0.2, 1.8, 2.5, 0.7, big, size=26, color=col, bold=True)
    text(s, x + 0.2, 2.6, 2.5, 0.4, t, size=13, color=INK, bold=True)
    text(s, x + 0.2, 3.0, 2.5, 0.9, sub, size=10.5, color=GREY)
big_panel(s, 0.55, 4.35, 12.23, 1.9, CARD)
text(s, 0.85, 4.6, 11.6, 0.4, "REAL EXAMPLE — restaurants running ads on Zomato and Swiggy",
     size=13, color=NAVY, bold=True)
text(s, 0.85, 5.05, 11.6, 1.1,
     "The best one spent ₹1.6L on ads in a month — and 7 in 10 orders came from those ads. The dashboard "
     "says ₹8.1 back for every ₹1, but that number counts all orders, even ones that would have come anyway. "
     "Real risk: if ads pause, revenue pauses too. Our plan: measure what ads truly add, build repeat "
     "customers, and lift the rating.", size=12.5, color=INK)
text(s, 0.55, 6.62, 12.2, 0.3,
     "Reported returns can lie. We test with real money — cut ads 20% for 2 weeks and watch what survives.",
     size=10.5, color=GREY)

# ================================================================ 08 THE MENU MACHINE
s = slide()
header(s, "THE MENU MACHINE", "Menu and offers: where profit quietly gets lost", "08")
menu = [
    (0.55, ORANGE, "ORDER SIZE", "Combos and add-ons take ₹200–300 orders to ₹350+. Small orders are bad — the platform keeps ~33% of them."),
    (6.75, RED, "OFFER DISCIPLINE", "A 10% discount is like paying 31% extra commission. We grow with platform-funded offers, never your profit."),
    (0.55, NAVY, "MENU CLEAN-UP", "Keep your best dishes on top, remove slow ones, test prices. The menu becomes a money-maker, not just a list."),
    (6.75, GREEN, "PACKING & FEES", "One clear packing charge for all items — no more ₹7 to ₹40 on the same menu that waste money and confuse tax."),
]
for x, col, t, d in menu:
    strip_card(s, x, 1.5 if t != "MENU CLEAN-UP" else 4.0, 6.03, 2.3, col)
    text(s, x + 0.3, 1.75 if t != "MENU CLEAN-UP" else 4.25, 5.4, 0.5, t, size=14, color=col, bold=True)
    text(s, x + 0.3, 2.35 if t != "MENU CLEAN-UP" else 4.85, 5.4, 1.35, d, size=12, color=INK)
text(s, 0.55, 6.62, 12.2, 0.3,
     "Order size up ₹50 on 1,500 orders = +₹75,000 more every month, same kitchen.", size=10.5, color=GREY)

# ================================================================ 09 THE DELIVERABLES
s = slide()
header(s, "THE DELIVERABLES", "What you get from us", "09")
dels = [
    (0.55, NAVY, "WEEKLY\nREVENUE REPORT", "Summary, what went right, what went wrong, next 3 steps. Numbers on every page."),
    (3.7, ORANGE, "HEALTH SCORE\nDASHBOARD", "Your health dials every week. Watch red turn green."),
    (6.85, GREEN, "WEEKLY\nACTION PLAN", "One fix done every week, with the ₹ it added."),
    (10.0, RED, "MONTHLY\nREVIEW", "What we did, what it earned, what's next."),
]
for x, col, t, d in dels:
    strip_card(s, x, 1.6, 2.9, 2.9, col)
    text(s, x + 0.2, 1.95, 2.5, 1.0, t, size=14, color=col, bold=True)
    text(s, x + 0.2, 3.0, 2.5, 1.3, d, size=10, color=INK)
big_panel(s, 0.55, 4.7, 12.23, 1.55, CARD)
text(s, 0.85, 4.9, 11.6, 0.4, "THE WEEKLY REPORT — 4 parts, all numbers", size=13, color=NAVY, bold=True)
parts = [
    (0.85, "SUMMARY", "sales, money received, wastage, order size — vs last week"),
    (3.9, "WHAT WENT RIGHT", "the good things, and the ₹ they brought"),
    (6.95, "WHAT WENT WRONG", "the problems, and the ₹ they cost"),
    (10.0, "NEXT STEPS", "3 clear actions for this week"),
]
for x, t, d in parts:
    text(s, x, 5.35, 2.9, 0.4, t, size=12, color=ORANGE, bold=True)
    text(s, x, 5.7, 2.9, 0.5, d, size=10, color=GREY)
big_panel(s, 0.55, 6.45, 12.23, 0.62, NAVY)
text(s, 0.85, 6.6, 11.7, 0.4, "Every page compares this week with last week. You read it in 3 minutes.",
     size=12.5, color=WHITE, bold=True)
text(s, 0.55, 6.62, 12.2, 0.3, "", size=10.5, color=GREY)  # spacer keeps footer rhythm

# ================================================================ 10 PROOF
s = slide()
header(s, "PROOF, NOT PROMISES", "What we found in two restaurants' real data", "10")
# --- Restaurant A
big_panel(s, 0.55, 1.45, 6.1, 5.05, CARD)
rect(s, 0.55, 1.45, 6.1, 0.09, ORANGE)
rect(s, 0.55, 1.45, 0.09, 5.05, ORANGE)
text(s, 0.85, 1.7, 5.6, 0.4, "RESTAURANT A — JULY 2026 · 3,749 ORDERS", size=12, color=NAVY, bold=True)
text(s, 0.85, 2.1, 5.6, 0.4,
     [[("HEALTH 86 ", {"bold": True, "color": GREEN}), ("· TRACK 2 — GROWTH", {"color": GREY})]],
     size=12.5)
a_findings = [
    (ORANGE, "Ads drive 7 in 10 orders", "₹1.6L spent in July — dashboard says 8.1× back, but that counts orders that would come anyway."),
    (NAVY, "Platform takes 20.7% of menu value", "commission + payment fee + tax on fees. Quiet, but it is the biggest cut."),
    (RED, "Rating 4.16 — below the 4.2 comfort line", "only a fraction of orders are rated; a fragile rating hurts search rank."),
    (GREEN, "Repeat customers untapped", "4 in 10 orders come back — a repeat offer is the cheapest growth left."),
]
y = 2.65
for col, t, d in a_findings:
    dot(s, 0.85, y + 0.04, col)
    text(s, 1.15, y, 4.9, 0.4, t, size=12.5, color=INK, bold=True)
    text(s, 1.15, y + 0.3, 4.9, 0.7, d, size=10.5, color=GREY)
    y += 1.05
# --- Restaurant B
big_panel(s, 6.8, 1.45, 6.1, 5.05, CARD)
rect(s, 6.8, 1.45, 6.1, 0.09, ORANGE)
rect(s, 6.8, 1.45, 0.09, 5.05, ORANGE)
text(s, 7.1, 1.7, 5.6, 0.4, "RESTAURANT B — JULY 2026 · 2,194 ORDERS", size=12, color=NAVY, bold=True)
text(s, 7.1, 2.1, 5.6, 0.4,
     [[("HEALTH 72 ", {"bold": True, "color": ORANGE}), ("· TRACK 2 — GROWTH", {"color": GREY})]],
     size=12.5)
b_findings = [
    (RED, "Nearly 9 in 10 orders come from ads", "87% dependency — pause the ads and revenue pauses with them."),
    (RED, "Orders fell month-on-month even with ads on", "the warning light: ads can't mask a shrinking base forever."),
    (NAVY, "Platform takes 29.7% of menu value", "the highest cut of the two — small orders and distance fees drive it."),
    (ORANGE, "Menu-to-cart below 30%", "people see the menu but don't add to cart — the menu page is the leak."),
]
y = 2.65
for col, t, d in b_findings:
    dot(s, 7.1, y + 0.04, col)
    text(s, 7.4, y, 4.9, 0.4, t, size=12.5, color=INK, bold=True)
    text(s, 7.4, y + 0.3, 4.9, 0.7, d, size=10.5, color=GREY)
    y += 1.05
text(s, 0.55, 6.62, 12.2, 0.3,
     "Real numbers from two real restaurants — this is the free check-up you get in 30 minutes.",
     size=10.5, color=GREY)

# ================================================================ 11 LET US PROVE IT
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY_D)
rect(s, 0, 0, 13.333, 0.12, ORANGE)
text(s, 12.1, 0.32, 0.9, 0.4, "11", size=11, color=GREY2, align=PP_ALIGN.RIGHT)
text(s, 0.9, 0.9, 11.5, 0.5, "LET US PROVE IT", size=14, color=ORANGE, bold=True)
text(s, 0.9, 1.3, 11.5, 1.0, "A free 30-minute check-up. No contract. No charges.",
     size=34, color=WHITE, bold=True)
steps = [
    ("1", "Send your last 3 reports", "Zomato or Swiggy — that's all we need"),
    ("2", "We study them overnight", "You get your Health Score, the ₹ that is leaking, and 3 fixes you can do this week"),
    ("3", "You decide", "Try us for 90 days. If your money doesn't grow, you walk away"),
]
x = 0.9
for num, t, d in steps:
    big_panel(s, x, 2.6, 3.6, 3.0, NAVY)
    text(s, x + 0.3, 2.9, 0.9, 0.8, num, size=36, color=GOLD, bold=True)
    text(s, x + 0.3, 3.75, 3.0, 0.8, t, size=14, color=WHITE, bold=True)
    text(s, x + 0.3, 4.5, 3.0, 1.0, d, size=11, color=BLUE1)
    x += 3.9
# tool strip — 3-point arrow diagram
big_panel(s, 0.9, 5.85, 11.5, 0.6, NAVY2)
arrow_parts = [
    "Your payout + funnel report",
    "Rules engine — 60+ checks per order",
    "Score, leaks & 3 fixes for this week",
]
x = 1.1
for i, part in enumerate(arrow_parts):
    text(s, x, 6.0, 3.3, 0.4, part, size=11, color=WHITE, bold=True)
    x += 3.45
    if i < 2:
        text(s, x - 0.25, 5.97, 0.4, 0.4, "→", size=14, color=GOLD, bold=True)
text(s, 0.9, 6.0 + 0.62, 11.5, 0.5,
     "The check-up is free because the numbers speak for themselves. We only earn when your money grows.",
     size=14, color=GOLD, bold=True)
text(s, 0.9, 6.9, 11.5, 0.4,
     "SCALE PLATES — Restaurant Revenue Intelligence   |   we look after the money side of your online restaurant",
     size=11, color=GREY2)

# ---------------------------------------------------------------- save
out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "pitch")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "Scale-Plates-Client-Pitch.pptx")
prs.save(out)
print("saved:", out)
print("slides:", len(list(prs.slides)))
